import streamlit as st
from pptx import Presentation
import io


def replace_text_in_shapes(shapes, old_text, new_text):
    """
    Finds and replaces text in all text-containing shapes.

    Args:
        shapes: A collection of shapes from a slide.
        old_text (str): The text to be replaced (e.g., "[TOPIC]").
        new_text (str): The new text to insert.
    """
    for shape in shapes:
        if shape.has_text_frame:
            # Check if the old text exists in the shape's text
            if old_text in shape.text:
                # Iterate through all paragraphs and runs to find and replace the text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)


def update_agenda_slide(slide, new_agenda_items):
    """
    Update agenda items inside GroupShape (SmartArt converted to shapes).
    """
    def update_shapes(shapes, new_items):
        for shape in shapes:
            if shape.shape_type == 6:  # GroupShape
                update_shapes(shape.shapes, new_items)

            elif shape.has_text_frame:
                text = shape.text.strip()
                for i, new_text in enumerate(new_items, start=1):
                    placeholder = f"[TITLE-{i}]"
                    if placeholder in text:
                        shape.text = new_text

    # Start recursive update
    update_shapes(slide.shapes, new_agenda_items)


def update_subtopic_slide(slide, replacements, new_images=None):
    """
    Update text placeholders and replace multiple pictures in the slide.

    :param slide: pptx slide object
    :param replacements: dict {placeholder -> new text}
    :param new_images: list of image paths to replace pictures in order of appearance
    """
    img_index = 0

    for shape in list(slide.shapes):  # make list to allow modifications
        # --- Replace text placeholders ---
        if shape.has_text_frame:
            for placeholder, new_text in replacements.items():
                if placeholder in shape.text:
                    print(f"Updating {placeholder} -> {new_text}")
                    shape.text = new_text

        # --- Replace pictures ---
        if new_images and shape.shape_type == 13:  # 13 = Picture
            if img_index < len(new_images):
                new_image_path = new_images[img_index]
                print(f"Replacing Picture {img_index+1} with {new_image_path}")
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(
                    shape._element)  # remove old picture
                slide.shapes.add_picture(
                    new_image_path, left, top, width, height)
                img_index += 1


# Set up the Streamlit page
st.set_page_config(page_title="PPT Generator", layout="centered")

st.title("Automated PPT Generator")
st.markdown(
    "Easily create your presentations from a template by filling out the form below.")

# Create a form to gather user input
with st.form(key="ppt_form"):
    st.subheader("Slide 1")
    topic_title = st.text_input(
        "Topic Title", help="Enter the main title for your presentation (e.g., 'Introduction to Python').", key="topic_title")
    subtopics = st.text_input(
        "Subtopics", help="Enter the subtopics, separated by , (one subtopic per line).", key="subtopics_area")
    st.subheader("Slide 2")
    agenta_title = st.text_input(
        'Agenda title', help="Enter the Agenda, separated by ,", key="agenda_title")

    # Collect data for each content slide dynamically
    content_slides_data = []
    for i in range(1, 14):  # Slides 3 to 14
        st.subheader(f"Slide {i + 2}")
        slide_subtopic = st.text_input(
            'Slide Topic', help="Enter the topic of the slide", key=f"slide_subtopic_{i}")
        syntax = st.text_area(
            'Syntax', help="Enter the syntax for the slide", key=f"syntax_{i}")
        explain = st.text_area(
            'Explanation', help="Enter the explanation for the slide", key=f"explain_{i}")
        image1 = st.file_uploader("Upload Implementation Image", type=[
                                  "png", "jpg", "jpeg"], key=f"image1_{i}")
        image2 = st.file_uploader("Upload Output Image", type=[
                                  "png", "jpg", "jpeg"], key=f"image2_{i}")

        content_slides_data.append({
            'slide_index': i+1,
            'subtopic': slide_subtopic,
            'syntax': syntax,
            'explain': explain,
            'images': [image1, image2]
        })

    # Slide 16
    st.subheader(f"Slide 16")
    last_summary_slide_text = st.text_area(
        "Last Slide Text", help="Enter the text Summary for last slide (eg: Introduction,While Loops,For Loops,Nested Loops,Loop Control Statements,Randomness in Loops,Applications,Common Mistakes,Best Practices,Summary).", key="last_slide")

    submit_button = st.form_submit_button(label="Generate and Download PPT")
# Process the form submission
if submit_button and topic_title and subtopics:
    try:
        # Load the presentation template
        prs = Presentation("Template.pptx")

        # Check if the template has the expected number of slides
        if len(prs.slides) != 16:
            st.error(
                f"Expected a template with exactly 15 slides, but found {len(prs.slides)}. Please check your template.")
            st.stop()

        # Slide 1
        slide1 = prs.slides[0]
        replace_text_in_shapes(slide1.shapes, "[TOPIC]", topic_title)
        replace_text_in_shapes(slide1.shapes, "[SUBTOPICS]", subtopics)

        # Slide 2
        slide2 = prs.slides[1]
        update_agenda_slide(slide2, agenta_title.split(','))

        # Slides 3 to 15
        # Update the content slides (looping through your template's content slides)
        for data in content_slides_data:
            slide_index = data['slide_index']

            content_slide = prs.slides[slide_index]

            replacements = {
                "[SUBTOPIC]": data['subtopic'],
                "[SUBTOPIC EXPLAINS]": data['subtopic'] + 'Explained',
                "[SYNTAX]": data['syntax'],
                "[EXPLAIN]": data['explain']
            }
            images_to_use = [img for img in data['images'] if img is not None]
            update_subtopic_slide(content_slide, replacements, images_to_use)

        # Slide 16
        summary_slide = prs.slides[15]
        update_agenda_slide(summary_slide, last_summary_slide_text.split(','))

        # Save the new presentation to a BytesIO object in memory
        output = io.BytesIO()
        prs.save(output)

        # Create a download button for the user
        st.success(
            "Your presentation has been generated! Click the button below to download.")
        st.download_button(
            label="Download PPTX File",
            data=output.getvalue(),
            file_name=f"{topic_title.replace(' ', '_').lower()}_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except FileNotFoundError:
        st.error(
            "Template file not found. Please ensure 'Template.pptx' is in the application directory.")

    except Exception as e:
        st.error(f"An error occurred: {e}")
